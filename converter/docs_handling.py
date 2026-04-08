"""
Document converters: txt, md, docx, pdf, pptx, xlsx.
"""
import logging
from pathlib import Path
from typing import Optional

from .base import BaseConverter
from .types import ConvertedContent, ConversionMetadata

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
    import pandas as pd
except ImportError:
    openpyxl = None
    pd = None

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None


logger = logging.getLogger(__name__)


class TextConverter(BaseConverter):
    """Convert .txt files to text."""

    def can_handle(self, file_path: str) -> bool:
        return str(file_path).lower().endswith(".txt")

    def convert(self, file_path: str) -> ConvertedContent:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
            text = self.normalize_text(text)
            metadata = self.create_metadata(file_path, "txt", "text")
            return ConvertedContent(text=text, metadata=metadata, success=True)
        except Exception as e:
            logger.error(f"Error converting {file_path}: {e}")
            return ConvertedContent(
                text="",
                metadata=self.create_metadata(file_path, "txt", "text"),
                success=False,
                error=str(e),
            )


class MarkdownConverter(BaseConverter):
    """Convert .md files to text."""

    def can_handle(self, file_path: str) -> bool:
        return str(file_path).lower().endswith(".md")

    def convert(self, file_path: str) -> ConvertedContent:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
            text = self.normalize_text(text)
            metadata = self.create_metadata(file_path, "md", "text")
            return ConvertedContent(text=text, metadata=metadata, success=True)
        except Exception as e:
            logger.error(f"Error converting {file_path}: {e}")
            return ConvertedContent(
                text="",
                metadata=self.create_metadata(file_path, "md", "text"),
                success=False,
                error=str(e),
            )


class DocxConverter(BaseConverter):
    """Convert .docx files to text."""

    def can_handle(self, file_path: str) -> bool:
        if not DocxDocument:
            return False
        return str(file_path).lower().endswith(".docx")

    def convert(self, file_path: str) -> ConvertedContent:
        if not DocxDocument:
            return ConvertedContent(
                text="",
                metadata=self.create_metadata(file_path, "docx", "text"),
                success=False,
                error="python-docx not installed",
            )
        try:
            doc = DocxDocument(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            text = "\n\n".join(paragraphs)
            text = self.normalize_text(text)
            metadata = self.create_metadata(file_path, "docx", "text")
            return ConvertedContent(text=text, metadata=metadata, success=True)
        except Exception as e:
            logger.error(f"Error converting {file_path}: {e}")
            return ConvertedContent(
                text="",
                metadata=self.create_metadata(file_path, "docx", "text"),
                success=False,
                error=str(e),
            )


class PdfConverter(BaseConverter):
    """Convert .pdf files to text using PyMuPDF or pdfplumber."""

    def can_handle(self, file_path: str) -> bool:
        if not (fitz or pdfplumber):
            return False
        return str(file_path).lower().endswith(".pdf")

    def convert(self, file_path: str) -> ConvertedContent:
        if not (fitz or pdfplumber):
            return ConvertedContent(
                text="",
                metadata=self.create_metadata(file_path, "pdf", "text"),
                success=False,
                error="PyMuPDF or pdfplumber not installed",
            )

        try:
            # Try pdfplumber first (better layout awareness)
            if pdfplumber:
                return self._convert_with_pdfplumber(file_path)
            else:
                return self._convert_with_fitz(file_path)
        except Exception as e:
            logger.error(f"Error converting {file_path}: {e}")
            return ConvertedContent(
                text="",
                metadata=self.create_metadata(file_path, "pdf", "text"),
                success=False,
                error=str(e),
            )

    def _convert_with_pdfplumber(self, file_path: str) -> ConvertedContent:
        """Convert PDF using pdfplumber."""
        texts_by_page = []
        with pdfplumber.open(file_path) as pdf:
            for page_idx, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text:
                    texts_by_page.append(text)

        full_text = "\n\n".join(texts_by_page)
        full_text = self.normalize_text(full_text)
        metadata = self.create_metadata(
            file_path, "pdf", "text", num_pages=len(pdf.pages)
        )
        return ConvertedContent(text=full_text, metadata=metadata, success=True)

    def _convert_with_fitz(self, file_path: str) -> ConvertedContent:
        """Convert PDF using PyMuPDF."""
        doc = fitz.open(file_path)
        texts_by_page = []
        for page_idx, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                texts_by_page.append(text)

        doc.close()
        full_text = "\n\n".join(texts_by_page)
        full_text = self.normalize_text(full_text)
        metadata = self.create_metadata(
            file_path, "pdf", "text", num_pages=len(doc)
        )
        return ConvertedContent(text=full_text, metadata=metadata, success=True)


class PptxConverter(BaseConverter):
    """Convert .pptx (PowerPoint) files to text."""

    def can_handle(self, file_path: str) -> bool:
        if not Presentation:
            return False
        return str(file_path).lower().endswith(".pptx")

    def convert(self, file_path: str) -> ConvertedContent:
        if not Presentation:
            return ConvertedContent(
                text="",
                metadata=self.create_metadata(file_path, "pptx", "text"),
                success=False,
                error="python-pptx not installed",
            )
        try:
            prs = Presentation(file_path)
            texts_by_slide = []
            for slide_idx, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)
                if slide_texts:
                    texts_by_slide.append(" ".join(slide_texts))

            text = "\n\n".join(texts_by_slide)
            text = self.normalize_text(text)
            metadata = self.create_metadata(
                file_path, "pptx", "text", num_slides=len(prs.slides)
            )
            return ConvertedContent(text=text, metadata=metadata, success=True)
        except Exception as e:
            logger.error(f"Error converting {file_path}: {e}")
            return ConvertedContent(
                text="",
                metadata=self.create_metadata(file_path, "pptx", "text"),
                success=False,
                error=str(e),
            )


class ExcelConverter(BaseConverter):
    """Convert .xlsx (Excel) files to text."""

    def can_handle(self, file_path: str) -> bool:
        if not (openpyxl or pd):
            return False
        return str(file_path).lower().endswith((".xlsx", ".xls"))

    def convert(self, file_path: str) -> ConvertedContent:
        if not (openpyxl or pd):
            return ConvertedContent(
                text="",
                metadata=self.create_metadata(file_path, "xlsx", "text"),
                success=False,
                error="openpyxl or pandas not installed",
            )
        try:
            # Use pandas for easier flattening
            if pd:
                return self._convert_with_pandas(file_path)
            else:
                return self._convert_with_openpyxl(file_path)
        except Exception as e:
            logger.error(f"Error converting {file_path}: {e}")
            return ConvertedContent(
                text="",
                metadata=self.create_metadata(file_path, "xlsx", "text"),
                success=False,
                error=str(e),
            )

    def _convert_with_pandas(self, file_path: str) -> ConvertedContent:
        """Convert Excel using pandas."""
        excel_file = pd.ExcelFile(file_path)
        texts_by_sheet = []

        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            sheet_text = f"Sheet: {sheet_name}\n"
            sheet_text += df.to_string()
            texts_by_sheet.append(sheet_text)

        text = "\n\n".join(texts_by_sheet)
        text = self.normalize_text(text)
        metadata = self.create_metadata(
            file_path,
            "xlsx",
            "text",
            num_sheets=len(excel_file.sheet_names),
            sheet_names=excel_file.sheet_names,
        )
        return ConvertedContent(text=text, metadata=metadata, success=True)

    def _convert_with_openpyxl(self, file_path: str) -> ConvertedContent:
        """Convert Excel using openpyxl (fallback)."""
        wb = openpyxl.load_workbook(file_path, data_only=True)
        texts_by_sheet = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_text = f"Sheet: {sheet_name}\n"
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_str = " | ".join(str(cell) if cell is not None else "" for cell in row)
                rows.append(row_str)
            sheet_text += "\n".join(rows)
            texts_by_sheet.append(sheet_text)

        wb.close()
        text = "\n\n".join(texts_by_sheet)
        text = self.normalize_text(text)
        metadata = self.create_metadata(
            file_path, "xlsx", "text", num_sheets=len(wb.sheetnames)
        )
        return ConvertedContent(text=text, metadata=metadata, success=True)
